"""
PPTX Parser Module

This module provides functionality to parse PPTX (Microsoft PowerPoint) files using:
- python-pptx: Extract text, images, shapes, and slides

The main function is parse_pptx(), which extracts content and metadata
from PPTX files and returns them in a structured format.
"""

import os
import logging
import tempfile
from typing import Dict, List, Any, Union

# Import utility functions
from .utils import safe_parse, create_result_dict

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Try to import the required library
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PPTX parsing will not work.")

def extract_text_from_shape(shape) -> str:
    """
    Extract text from a PowerPoint shape.
    
    Args:
        shape: A shape object from python-pptx
        
    Returns:
        str: Extracted text or empty string if no text
    """
    if hasattr(shape, "text") and shape.text:
        return shape.text.strip()
    return ""

def extract_text_from_slide(slide) -> str:
    """
    Extract all text from a PowerPoint slide.
    
    Args:
        slide: A slide object from python-pptx
        
    Returns:
        str: Extracted text
    """
    texts = []
    
    # Extract text from shapes
    for shape in slide.shapes:
        text = extract_text_from_shape(shape)
        if text:
            texts.append(text)
            
    # Extract text from tables
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text:
                        texts.append(cell.text.strip())
    
    return "\n".join(texts)

def extract_tables_from_slide(slide) -> List[List[List[str]]]:
    """
    Extract tables from a PowerPoint slide.
    
    Args:
        slide: A slide object from python-pptx
        
    Returns:
        List[List[List[str]]]: List of tables, where each table is a list of rows,
                              and each row is a list of cell values
    """
    tables = []
    
    # Extract tables from shapes
    for shape in slide.shapes:
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip() if cell.text else "")
                table_data.append(row_data)
            tables.append(table_data)
    
    return tables

def extract_images_from_slide(slide, presentation_file, slide_idx) -> List[Dict[str, Any]]:
    """
    Extract images from a PowerPoint slide and save them to temporary files.
    
    Args:
        slide: A slide object from python-pptx
        presentation_file: Original PPTX file path for naming
        slide_idx: Index of the slide
        
    Returns:
        List[Dict[str, Any]]: List of image info dictionaries
    """
    images = []
    shape_idx = 0
    
    for shape in slide.shapes:
        shape_idx += 1
        
        if hasattr(shape, "image"):
            try:
                # Create a unique filename for the image
                base_filename = os.path.basename(presentation_file)
                temp_dir = tempfile.gettempdir()
                image_filename = f"{base_filename}_slide{slide_idx}_shape{shape_idx}.{shape.image.ext}"
                image_path = os.path.join(temp_dir, image_filename)
                
                # Save the image
                with open(image_path, 'wb') as f:
                    f.write(shape.image.blob)
                
                # Add image info to results
                images.append({
                    "path": image_path,
                    "slide_index": slide_idx,
                    "shape_index": shape_idx,
                    "format": shape.image.ext,
                    "size": len(shape.image.blob)
                })
            except Exception as e:
                logger.error(f"Error extracting image: {str(e)}")
    
    return images

def extract_slide_content(slide, presentation_file, slide_idx) -> Dict[str, Any]:
    """
    Extract all content from a PowerPoint slide.
    
    Args:
        slide: A slide object from python-pptx
        presentation_file: Original PPTX file path for naming
        slide_idx: Index of the slide
        
    Returns:
        Dict[str, Any]: Slide content
    """
    content = {
        "index": slide_idx,
        "text": extract_text_from_slide(slide)
    }
    
    # Extract tables if any
    tables = extract_tables_from_slide(slide)
    if tables:
        content["tables"] = tables
    
    # Extract images if any
    images = extract_images_from_slide(slide, presentation_file, slide_idx)
    if images:
        content["images"] = images
    
    # Extract slide title if available
    title_shape = None
    for shape in slide.shapes:
        if shape.is_title:
            title_shape = shape
            break
    
    if title_shape:
        content["title"] = extract_text_from_shape(title_shape)
    
    # Get slide layout name
    content["layout"] = slide.slide_layout.name
    
    return content

def extract_metadata_from_pptx(presentation) -> Dict[str, Any]:
    """
    Extract metadata from a PowerPoint presentation.
    
    Args:
        presentation: A Presentation object from python-pptx
        
    Returns:
        Dict[str, Any]: Metadata dictionary
    """
    metadata = {}
    
    # Core properties
    if hasattr(presentation, "core_properties"):
        props = presentation.core_properties
        if props.author:
            metadata["author"] = props.author
        if props.title:
            metadata["title"] = props.title
        if props.subject:
            metadata["subject"] = props.subject
        if props.created:
            metadata["created"] = props.created.isoformat() if hasattr(props.created, "isoformat") else str(props.created)
        if props.modified:
            metadata["modified"] = props.modified.isoformat() if hasattr(props.modified, "isoformat") else str(props.modified)
        if props.comments:
            metadata["comments"] = props.comments
    
    # Presentation stats
    metadata["slide_count"] = len(presentation.slides)
    
    return metadata

def _parse_pptx_internal(file_path: str) -> Dict[str, Any]:
    """
    Internal function to parse a PPTX file.
    
    Args:
        file_path (str): Path to the PPTX file
        
    Returns:
        Dict[str, Any]: Parsed content
    """
    if not PPTX_AVAILABLE:
        return create_result_dict(error="python-pptx is not available for PPTX parsing.")
    
    try:
        # Load the presentation
        presentation = Presentation(file_path)
        
        # Extract metadata
        metadata = extract_metadata_from_pptx(presentation)
        metadata["filename"] = os.path.basename(file_path)
        metadata["filesize"] = os.path.getsize(file_path)
        
        # Extract slides content
        slides = []
        full_text = []
        
        for idx, slide in enumerate(presentation.slides):
            slide_content = extract_slide_content(slide, file_path, idx)
            slides.append(slide_content)
            
            # Add slide text to full text
            if "text" in slide_content and slide_content["text"]:
                if "title" in slide_content:
                    full_text.append(f"Slide {idx+1} - {slide_content['title']}")
                else:
                    full_text.append(f"Slide {idx+1}")
                full_text.append(slide_content["text"])
                full_text.append("")  # Empty line between slides
        
        # Prepare result
        result = {
            "metadata": metadata,
            "slides": slides,
            "text": "\n".join(full_text)
        }
        
        return result
    except Exception as e:
        return create_result_dict(error=f"Error parsing PPTX file: {str(e)}")

def parse_pptx(file_path: str) -> Dict[str, Any]:
    """
    Parse a PPTX file and extract text, tables, images, and metadata.
    
    Args:
        file_path (str): Path to the PPTX file
        
    Returns:
        Dict[str, Any]: A dictionary containing:
            - text (str): Full extracted text from all slides
            - slides (List[Dict]): Content of each slide
            - metadata (Dict[str, Any]): File metadata
            - error (str): Error message (if any)
    """
    return safe_parse(_parse_pptx_internal, file_path)

# Example usage
if __name__ == "__main__":
    import sys
    import json
    
    if len(sys.argv) < 2:
        print("Usage: python pptx_parser.py <pptx_file_path>")
        sys.exit(1)
        
    file_path = sys.argv[1]
    result = parse_pptx(file_path)
    
    if "error" in result:
        print(f"Error: {result['error']}")
    else:
        print(f"Presentation: {result['metadata'].get('title', 'Untitled')}")
        print(f"Slides: {result['metadata']['slide_count']}")
        print(f"Text (excerpt): {result['text'][:300]}...")
        
        # Count total tables and images
        table_count = 0
        image_count = 0
        for slide in result["slides"]:
            if "tables" in slide:
                table_count += len(slide["tables"])
            if "images" in slide:
                image_count += len(slide["images"])
        
        print(f"Tables found: {table_count}")
        print(f"Images found: {image_count}")
        print(f"Full result saved to {file_path}.json")
        
        # Save full result to JSON file
        with open(f"{file_path}.json", "w") as f:
            json.dump(result, f, indent=2) 